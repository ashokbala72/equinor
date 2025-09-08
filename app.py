import os
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
from dotenv import load_dotenv
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Inches

# ---------------------------
# Load secrets from .env
# ---------------------------
load_dotenv()

AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview")
DEPLOYMENT_NAME = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o-raj")

client = AzureOpenAI(
    api_key=AZURE_OPENAI_API_KEY,
    api_version=AZURE_OPENAI_API_VERSION,
    azure_endpoint=AZURE_OPENAI_ENDPOINT
)

def get_genai_findings(data_batch: str, focus_area: str) -> str:
    prompt = f"""
    You are an ESG/energy sustainability expert.
    Analyze the following {focus_area} dataset and produce:
    1. Key Findings (summary of trends)
    2. Recommendations (clear action items with rationale)

    Dataset:
    {data_batch}
    """
    try:
        response = client.chat.completions.create(
            model=DEPLOYMENT_NAME,
            messages=[
                {"role": "system", "content": "You are an expert ESG advisor for an energy company."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"⚠️ GenAI analysis failed: {e}"

# ---------------------------
# Streamlit Config
# ---------------------------
st.set_page_config(page_title="Equinor Sustainability Dashboard", layout="wide")

st.markdown("""
    <style>
        body, .stApp { background-color: #ffffff; color: #000000; }
        h1, h2, h3, h4, h5 { color: #000000; }
        .stTabs [role="tab"] { background-color: #f0f0f0; color: #000000; padding: 8px; }
        .stTabs [role="tab"]:hover { background-color: #cccccc; }
        .banner { color:black; padding:6px; background-color:#ddd; border-radius:4px; margin-bottom:10px; }
        table, th, td { font-weight: bold !important; color: black !important; }
    </style>
""", unsafe_allow_html=True)

st.title("🌍 Equinor Sustainability Dashboard")

# ---------------------------
# Example DataFrames (simulated)
# ---------------------------
df_energy = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Oil": [50, 48, 46, 44, 42],
    "Gas": [30, 32, 33, 34, 35],
    "Renewables": [20, 20, 21, 22, 23]
})

df_emissions = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Scope 1": [10.5, 9.8, 9.3, 8.7, 8.3],
    "Scope 2": [2.1, 2.0, 1.9, 1.7, 1.6],
    "Scope 3": [45.0, 44.5, 44.0, 43.2, 42.5]
})

df_env = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Spills": [15, 12, 10, 8, 7],
    "Waste (kt)": [200, 210, 205, 198, 190]
})

df_safety = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "TRIF": [2.5, 2.2, 2.0, 1.8, 1.6],
    "SIF": [0.5, 0.4, 0.4, 0.3, 0.3]
})

df_workforce = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Employees": [20000, 20500, 21000, 21500, 22000],
    "Female %": [28, 30, 31, 33, 35],
    "Attrition %": [8, 7.5, 7.8, 8.2, 8.0]
})

df_socio = pd.DataFrame({
    "Country": ["Norway", "Brazil", "UK"],
    "Taxes Paid (MUSD)": [1200, 800, 600],
    "Local Procurement (MUSD)": [500, 400, 350]
})

df_hr = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Supplier Audits": [120, 135, 140, 150, 160],
    "Grievances": [5, 8, 6, 7, 9]
})

df_gov = pd.DataFrame({
    "Year": [2020, 2021, 2022, 2023, 2024],
    "Board Independence %": [75, 77, 78, 80, 82],
    "Anti-Corruption Training %": [92, 94, 95, 96, 97]
})

df_targets = pd.DataFrame({
    "Target": ["CO2 Reduction by 2030", "Renewables Capacity 2030"],
    "Planned": [50, 12],
    "Actual": [38, 9]
})

# ---------------------------
# Tabs
# ---------------------------
tabs = st.tabs([
    "⚡ Energy Mix", "🌍 Climate & Emissions", "💧 Environment",
    "⛑️ Safety", "👥 Workforce", "💰 Socio-Economic",
    "🤝 Human Rights", "🏛️ Governance", "🎯 ESG Targets", "📝 Reporting"
])

report_content = {}

# ---------------------------
# Chart helpers (2D only)
# ---------------------------
DARK_TEAL = "#005f73"
GREEN = "#94d2bd"
ORANGE = "#ee9b00"

def plot_bar(df, categories, ylabel, colors, xlabel="Year", unit=""):
    fig, ax = plt.subplots(figsize=(8,5))
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")

    x = np.arange(len(df[xlabel]))
    width = 0.8 / len(categories)
    for i, col in enumerate(categories):
        b = ax.bar(x + (i - (len(categories)-1)/2) * width, df[col], width,
                   label=col, color=colors[i])
        ax.bar_label(b, padding=2, color="black", fontsize=10, fmt="%.1f")
    ax.set_xticks(x)
    ax.set_xticklabels(df[xlabel], color="black")
    ax.set_xlabel(f"{xlabel}", color="black", fontsize=12)
    ax.set_ylabel(f"{ylabel} {unit}", color="black", fontsize=12)
    ax.set_title(ylabel, color="black", fontsize=14)
    ax.tick_params(axis="x", colors="black")
    ax.tick_params(axis="y", colors="black")
    ax.legend(facecolor="#ffffff", labelcolor="black")
    return fig

def plot_line(df, column, ylabel, xlabel="Year", unit=""):
    fig, ax = plt.subplots(figsize=(8,5))
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")

    ax.plot(df[xlabel], df[column], marker="o", color=DARK_TEAL, linewidth=2, label=column)
    for i, val in enumerate(df[column]):
        ax.text(df[xlabel][i], val, f"{val:.1f}{unit}", ha="center", va="bottom",
                color="black", fontsize=10)

    ax.set_xlabel(f"{xlabel}", color="black", fontsize=12)
    ax.set_ylabel(f"{ylabel} {unit}", color="black", fontsize=12)
    ax.set_title(ylabel, color="black", fontsize=14)
    ax.tick_params(axis="x", colors="black")
    ax.tick_params(axis="y", colors="black")
    ax.legend(facecolor="#ffffff", labelcolor="black")
    return fig

def plot_bar_categorical(df, categories, xlabel, colors, ylabel, unit=""):
    fig, ax = plt.subplots(figsize=(8,5))
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")

    x = np.arange(len(df[xlabel]))
    width = 0.35
    bars1 = ax.bar(x - width/2, df[categories[0]], width, label=categories[0], color=colors[0])
    bars2 = ax.bar(x + width/2, df[categories[1]], width, label=categories[1], color=colors[1])
    ax.bar_label(bars1, padding=2, color="black", fontsize=10, fmt="%.1f")
    ax.bar_label(bars2, padding=2, color="black", fontsize=10, fmt="%.1f")
    ax.set_xticks(x)
    ax.set_xticklabels(df[xlabel], color="black")
    ax.set_xlabel(xlabel, color="black", fontsize=12)
    ax.set_ylabel(f"{ylabel} {unit}", color="black", fontsize=12)
    ax.set_title(ylabel, color="black", fontsize=14)
    ax.tick_params(axis="y", colors="black")
    ax.legend(facecolor="#ffffff", labelcolor="black")
    return fig

# ---------------------------
# Section renderer
# ---------------------------
def render_section(name, df, plot_func, source_info, focus_area):
    st.header(name)
    st.markdown(f'<div class="banner">🔵 Data is Simulated (should come from {source_info})</div>', unsafe_allow_html=True)
    st.dataframe(df, use_container_width=True)

    fig = plot_func(df)
    st.pyplot(fig)

    st.subheader("🔎 GenAI Findings & Recommendations")
    with st.spinner("Generating AI insights..."):
        findings = get_genai_findings(df.to_json(orient="records"), focus_area)
    st.markdown(findings)

    report_content[focus_area] = (fig, findings)

# ---------------------------
# Tab Content
# ---------------------------
with tabs[0]:
    render_section("⚡ Energy Mix", df_energy,
                   lambda df: plot_bar(df, ["Oil","Gas","Renewables"], "Energy Share", [DARK_TEAL, GREEN, ORANGE], xlabel="Year", unit="%"),
                   "Production Reporting System & Grid Operator Feeds",
                   "Energy Mix")

with tabs[1]:
    render_section("🌍 Climate & Emissions", df_emissions,
                   lambda df: plot_bar(df, ["Scope 1","Scope 2","Scope 3"], "Emissions", [DARK_TEAL, GREEN, ORANGE], xlabel="Year", unit="Mt CO₂e"),
                   "PI System / SCADA Sensors / GHG Monitoring Systems",
                   "Emissions")

with tabs[2]:
    render_section("💧 Environment", df_env,
                   lambda df: plot_line(df, "Waste (kt)", "Waste", xlabel="Year", unit="kt"),
                   "Environmental Monitoring System & LIMS",
                   "Environment")

with tabs[3]:
    render_section("⛑️ Safety", df_safety,
                   lambda df: plot_line(df, "TRIF", "TRIF", xlabel="Year", unit=" per million hours"),
                   "Synergi Life & Operational Risk Management Systems",
                   "Safety")

with tabs[4]:
    render_section("👥 Workforce", df_workforce,
                   lambda df: plot_bar(df, ["Employees","Female %","Attrition %"], "Workforce Metrics", [DARK_TEAL, GREEN, ORANGE], xlabel="Year"),
                   "SAP SuccessFactors HRIS / Payroll",
                   "Workforce")

with tabs[5]:
    render_section("💰 Socio-Economic Impact", df_socio,
                   lambda df: plot_bar_categorical(df, ["Taxes Paid (MUSD)","Local Procurement (MUSD)"], "Country", [DARK_TEAL, GREEN], "Value", unit="MUSD"),
                   "SAP ERP / Finance Systems",
                   "Socio-Economic")

with tabs[6]:
    render_section("🤝 Human Rights", df_hr,
                   lambda df: plot_bar(df, ["Supplier Audits","Grievances"], "Human Rights", [DARK_TEAL, GREEN], xlabel="Year", unit="count"),
                   "Supplier Risk & Grievance Management Systems",
                   "Human Rights")

with tabs[7]:
    render_section("🏛️ Governance", df_gov,
                   lambda df: plot_bar(df, ["Board Independence %","Anti-Corruption Training %"], "Governance", [DARK_TEAL, GREEN], xlabel="Year", unit="%"),
                   "SAP GRC & Board Governance Portal",
                   "Governance")

with tabs[8]:
    render_section("🎯 ESG Targets", df_targets,
                   lambda df: plot_bar(df, ["Planned","Actual"], "ESG Targets", [DARK_TEAL, ORANGE], xlabel="Target", unit="% or GW"),
                   "Sustainability KPI Tracker & Annual Report Database",
                   "ESG Targets")

with tabs[9]:
    st.header("📝 Reporting Tab")
    for section, (fig, findings) in report_content.items():
        st.subheader(section)
        st.pyplot(fig)
        st.markdown(findings)

    # PPT Export only
    prs = Presentation()
    for section, (fig, findings) in report_content.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = section
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = textbox.text_frame
        tf.text = findings
    pptx_output = BytesIO()
    prs.save(pptx_output)
    pptx_output.seek(0)

    st.download_button("📥 Download as PPT", data=pptx_output,
                       file_name="Sustainability_Report.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
