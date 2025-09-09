import os
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from openai import AzureOpenAI

# ---------------------------
# Load Azure Config
# ---------------------------
load_dotenv(override=True)
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview")
DEPLOYMENT_NAME = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o-raj")

client = AzureOpenAI(
    api_key=AZURE_OPENAI_API_KEY,
    api_version=AZURE_OPENAI_API_VERSION,
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
)

# ---------------------------
# GenAI Helper
# ---------------------------
def get_genai_findings(data_batch: str, focus_area: str) -> str:
    prompt = f"""
    You are an ESG/energy sustainability expert.
    Analyze the following {focus_area} dataset and produce:
    1. Key Findings
    2. Recommendations

    Dataset:
    {data_batch}
    """
    try:
        resp = client.chat.completions.create(
            model=DEPLOYMENT_NAME,
            messages=[
                {"role": "system", "content": "You are an ESG advisor."},
                {"role": "user", "content": prompt},
            ],
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"⚠️ GenAI failed: {e}"

# ---------------------------
# Equinor Colors
# ---------------------------
DARK_TEAL = "#005f73"
GREEN = "#94d2bd"
ORANGE = "#ee9b00"

# ---------------------------
# Chart Helpers
# ---------------------------
def plot_bar(df, categories, ylabel, colors, xlabel="Year", unit=""):
    fig, ax = plt.subplots(figsize=(7, 4))
    x = np.arange(len(df[xlabel]))
    width = 0.8 / len(categories)
    for i, col in enumerate(categories):
        b = ax.bar(
            x + (i - (len(categories)-1)/2) * width,
            df[col],
            width,
            label=col,
            color=colors[i]
        )
        ax.bar_label(b, padding=2, color="black", fontsize=9, fmt="%.1f")
    ax.set_xticks(x)
    ax.set_xticklabels(df[xlabel])
    ax.set_xlabel(xlabel)
    ax.set_ylabel(f"{ylabel} {unit}")
    ax.legend()
    return fig

def plot_line(df, cols, ylabel, xlabel="Year", unit=""):
    fig, ax = plt.subplots(figsize=(7, 4))
    for col in cols:
        ax.plot(df[xlabel], df[col], marker="o", linewidth=2, label=col)
        for i, val in enumerate(df[col]):
            ax.text(df[xlabel][i], val, f"{val:.1f}{unit}", ha="center", va="bottom", fontsize=8)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(f"{ylabel} {unit}")
    ax.set_title(ylabel)
    ax.legend()
    return fig

# ---------------------------
# Scenario Data Functions
# ---------------------------
def get_energy_mix(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Oil": [50, 48, 45, 42] if scenario=="Improved" else [50, 52, 55, 57],
        "Gas": [30, 32, 34, 36] if scenario=="Improved" else [30, 31, 32, 33],
        "Renewables": [20, 20, 21, 22] if scenario=="Improved" else [20, 17, 13, 10],
    })

def get_emissions(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Scope 1": [10.5, 9.8, 9.0, 8.5] if scenario=="Improved" else [10.5, 11.0, 12.0, 13.0],
        "Scope 2": [2.1, 2.0, 1.9, 1.8] if scenario=="Improved" else [2.1, 2.3, 2.5, 2.7],
        "Scope 3": [45, 44, 43, 42] if scenario=="Improved" else [45, 48, 50, 55],
    })

def get_environment(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Spills": [5, 4, 3, 2] if scenario=="Improved" else [5, 6, 7, 8],
        "Waste (kt)": [120, 110, 105, 100] if scenario=="Improved" else [120, 130, 140, 150],
    })

def get_safety(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "TRIF": [3.5, 3.0, 2.5, 2.0] if scenario=="Improved" else [3.5, 3.8, 4.2, 4.5],
        "SIF": [1.0, 0.9, 0.8, 0.7] if scenario=="Improved" else [1.0, 1.2, 1.4, 1.6],
    })

def get_workforce(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Employees": [20000, 20500, 21000, 21500] if scenario=="Improved" else [20000, 19800, 19500, 19000],
        "Diversity %": [28, 30, 33, 35] if scenario=="Improved" else [28, 27, 26, 25],
        "Attrition %": [8, 7.5, 7, 6.5] if scenario=="Improved" else [8, 9, 10, 11],
    })

def get_socio(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Taxes Paid (MUSD)": [2000, 2200, 2400, 2600] if scenario=="Improved" else [2000, 1800, 1600, 1500],
        "Local Procurement (MUSD)": [800, 850, 900, 950] if scenario=="Improved" else [800, 750, 700, 650],
    })

def get_human_rights(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Audits": [10, 12, 14, 15] if scenario=="Improved" else [10, 9, 8, 7],
        "Grievances": [5, 4, 3, 2] if scenario=="Improved" else [5, 6, 8, 9],
    })

def get_governance(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Board Independence %": [60, 65, 70, 72] if scenario=="Improved" else [60, 58, 55, 50],
        "Anti-Corruption Trainings": [500, 600, 650, 700] if scenario=="Improved" else [500, 450, 400, 350],
    })

def get_targets(scenario):
    return pd.DataFrame({
        "Year": [2020, 2021, 2022, 2023],
        "Planned": [100, 90, 80, 70],
        # Improved should show exceeding or meeting targets
        "Actual": [95, 92, 100, 110] if scenario=="Improved" else [95, 85, 75, 65],
    })


# ---------------------------
# Section Renderer
# ---------------------------
def render_section(name, scenario, data_func, plot_func, source_info, focus_area):
    st.subheader(f"{name} ({scenario} Scenario)")
    st.caption(f"🔵 Data is Simulated (should come from {source_info})")

    df = data_func(scenario)
    st.dataframe(df, use_container_width=True)

    fig = plot_func(df)
    st.pyplot(fig)

    with st.spinner("⏳ Generating GenAI recommendations..."):
        findings = get_genai_findings(df.to_csv(index=False), focus_area)

    st.markdown("### 🤖 GenAI Findings & Recommendations")
    st.write(findings)

    return (fig, findings)

# ---------------------------
# Streamlit UI
# ---------------------------
st.set_page_config(page_title="Equinor Dashboard", layout="wide")
st.title("🌍 Equinor Sustainability Dashboard")

tabs = st.tabs([
    "⚡ Energy Mix",
    "🌍 Emissions",
    "💧 Environment",
    "⛑️ Safety",
    "👥 Workforce",
    "💰 Socio-Economic",
    "🤝 Human Rights",
    "🏛️ Governance",
    "🎯 ESG Targets",
    "📝 Reporting"
])

report_content = {}

# ---------------------------
# Tabs 0–8
# ---------------------------
with tabs[0]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="energy_mix")
    fig, findings = render_section("⚡ Energy Mix", scenario, get_energy_mix,
                                   lambda d: plot_bar(d, ["Oil","Gas","Renewables"], "Energy Share", [DARK_TEAL, GREEN, ORANGE]),
                                   "Production Reporting System & Grid Feeds", "Energy Mix")
    report_content["Energy Mix"] = (fig, findings, scenario)

with tabs[1]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="emissions")
    fig, findings = render_section("🌍 Emissions", scenario, get_emissions,
                                   lambda d: plot_bar(d, ["Scope 1","Scope 2","Scope 3"], "Emissions", [DARK_TEAL, GREEN, ORANGE], unit="Mt CO₂e"),
                                   "SCADA & GHG Monitoring", "Emissions")
    report_content["Emissions"] = (fig, findings, scenario)

with tabs[2]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="environment")
    fig, findings = render_section(
        "💧 Environment",
        scenario,
        get_environment,
        lambda d: plot_line(d, ["Spills", "Waste (kt)"], "Environment KPIs", xlabel="Year"),
        "Environmental Monitoring Systems",
        "Environment"
    )
    report_content["Environment"] = (fig, findings, scenario)


with tabs[3]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="safety")
    fig, findings = render_section(
        "⛑️ Safety",
        scenario,
        get_safety,
        lambda d: plot_line(d, ["TRIF", "SIF"], "Safety KPIs", xlabel="Year"),
        "Safety Management Systems",
        "Safety"
    )
    report_content["Safety"] = (fig, findings, scenario)


with tabs[4]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="workforce")
    fig, findings = render_section("👥 Workforce", scenario, get_workforce,
                                   lambda d: plot_bar(d, ["Employees","Diversity %","Attrition %"], "Workforce", [DARK_TEAL, GREEN, ORANGE]),
                                   "HRIS", "Workforce")
    report_content["Workforce"] = (fig, findings, scenario)

with tabs[5]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="socio")
    fig, findings = render_section("💰 Socio-Economic", scenario, get_socio,
                                   lambda d: plot_bar(d, ["Taxes Paid (MUSD)","Local Procurement (MUSD)"], "Socio-Economic", [DARK_TEAL, ORANGE]),
                                   "Finance Systems", "Socio-Economic")
    report_content["Socio-Economic"] = (fig, findings, scenario)

with tabs[6]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="human_rights")
    fig, findings = render_section("🤝 Human Rights", scenario, get_human_rights,
                                   lambda d: plot_bar(d, ["Audits","Grievances"], "Human Rights", [DARK_TEAL, ORANGE]),
                                   "Supplier Audits", "Human Rights")
    report_content["Human Rights"] = (fig, findings, scenario)

with tabs[7]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="governance")
    fig, findings = render_section("🏛️ Governance", scenario, get_governance,
                                   lambda d: plot_bar(d, ["Board Independence %","Anti-Corruption Trainings"], "Governance", [DARK_TEAL, GREEN]),
                                   "GRC Systems", "Governance")
    report_content["Governance"] = (fig, findings, scenario)

with tabs[8]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="targets")
    fig, findings = render_section("🎯 ESG Targets", scenario, get_targets,
                                   lambda d: plot_bar(d, ["Planned","Actual"], "Targets", [DARK_TEAL, ORANGE]),
                                   "KPI Trackers", "ESG Targets")
    report_content["ESG Targets"] = (fig, findings, scenario)

# ---------------------------
# Reporting Tab
# ---------------------------
with tabs[9]:
    scenario = st.radio("Scenario", ["Improved", "Deteriorated"], horizontal=True, key="reporting")
    st.header(f"📝 Reporting ({scenario} Scenario)")

    prs = Presentation()
    for section, (fig, findings, scn) in report_content.items():
        if scn == scenario:
            st.subheader(section)
            st.pyplot(fig)
            st.markdown(findings)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"{section} ({scenario})"
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            textbox.text_frame.text = findings

    pptx_output = BytesIO()
    prs.save(pptx_output)
    pptx_output.seek(0)

    st.download_button(
        "📥 Download PPT",
        data=pptx_output,
        file_name=f"Sustainability_Report_{scenario}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
